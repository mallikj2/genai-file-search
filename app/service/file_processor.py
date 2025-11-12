import os
from typing import List, Dict
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd
import json
import xml.etree.ElementTree as ET
from loguru import logger

class FileProcessor:
    """Process different file types and extract text"""
    
    def process_file(self, file_path: str, file_type: str) -> List[Dict]:
        """
        Process file and return documents with text and metadata
        
        Args:
            file_path: Path to the file
            file_type: File extension
        
        Returns:
            List of document dictionaries
        """
        try:
            if file_type == ".pdf":
                return self._process_pdf(file_path)
            elif file_type in [".doc", ".docx"]:
                return self._process_docx(file_path)
            elif file_type in [".xls", ".xlsx"]:
                return self._process_excel(file_path)
            elif file_type in [".ppt", ".pptx"]:
                return self._process_pptx(file_path)
            elif file_type == ".txt":
                return self._process_txt(file_path)
            elif file_type == ".csv":
                return self._process_csv(file_path)
            elif file_type in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
                return self._process_image(file_path)
            elif file_type == ".json":
                return self._process_json(file_path)
            elif file_type == ".xml":
                return self._process_xml(file_path)
            elif file_type == ".sql":
                return self._process_sql(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            raise
    
    def _process_pdf(self, file_path: str) -> List[Dict]:
        """Extract text from PDF"""
        documents = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    documents.append({
                        "text": text,
                        "metadata": {"page_number": page_num + 1}
                    })
        return documents
    
    def _process_docx(self, file_path: str) -> List[Dict]:
        """Extract text from DOCX"""
        doc = Document(file_path)
        text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return [{"text": text, "metadata": {}}]
    
    def _process_excel(self, file_path: str) -> List[Dict]:
        """Extract text from Excel"""
        documents = []
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                documents.append({
                    "text": "\n".join(rows),
                    "metadata": {"sheet_name": sheet_name}
                })
        return documents
    
    def _process_pptx(self, file_path: str) -> List[Dict]:
        """Extract text from PowerPoint"""
        documents = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            if text_parts:
                documents.append({
                    "text": "\n".join(text_parts),
                    "metadata": {"slide_number": slide_num + 1}
                })
        return documents
    
    def _process_txt(self, file_path: str) -> List[Dict]:
        """Extract text from TXT"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
        return [{"text": text, "metadata": {}}]
    
    def _process_csv(self, file_path: str) -> List[Dict]:
        """Extract text from CSV"""
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
        return [{"text": text, "metadata": {"rows": len(df), "columns": len(df.columns)}}]
    
    def _process_image(self, file_path: str) -> List[Dict]:
        """Extract text from image using OCR"""
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return [{"text": text, "metadata": {}}]
    
    def _process_json(self, file_path: str) -> List[Dict]:
        """Extract text from JSON"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        text = json.dumps(data, indent=2)
        return [{"text": text, "metadata": {}}]
    
    def _process_xml(self, file_path: str) -> List[Dict]:
        """Extract text from XML"""
        tree = ET.parse(file_path)
        root = tree.getroot()
        text = ET.tostring(root, encoding='unicode', method='text')
        return [{"text": text, "metadata": {}}]
    
    def _process_sql(self, file_path: str) -> List[Dict]:
        """Extract text from SQL"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
        return [{"text": text, "metadata": {}}]

file_processor = FileProcessor()
